#!/usr/bin/env python3
"""Generate FlowPOS Hackathon Presentation (PPTX) — bold monochrome style."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Helpers ──────────────────────────────────────────────────────────────────

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0x99)
GRAY_DARK = RGBColor(0x77, 0x77, 0x77)
GRAY_LIGHT = RGBColor(0xBB, 0xBB, 0xBB)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill=None, outline=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if outline:
        shape.line.color.rgb = outline
        shape.line.width = Pt(1)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, font_name="Calibri", alignment=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=18,
                  color=WHITE, bold=False, font_name="Calibri",
                  alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add text box with multiple paragraphs (from list of (text, size, color, bold) tuples)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            txt, sz, clr, bld = line
        else:
            txt, sz, clr, bld = line, font_size, color, bold
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        p.space_before = Pt(0)
        p.space_after = Pt(4)
    return txBox


def dark_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLACK
    return slide


def light_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


# ── Build Presentation ───────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

PAD_L = Inches(1.2)
PAD_T = Inches(0.8)
CONTENT_W = Inches(10.9)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_text(s, Inches(0), Inches(2.0), SLIDE_W, Inches(0.4),
         "DIH x AUC HACKATHON", 13, GRAY_DARK, False, "Consolas", PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(2.5), SLIDE_W, Inches(1.6),
         "FlowPOS", 80, WHITE, True, "Calibri", PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.2), SLIDE_W, Inches(1.0),
         "AI-Powered Demand Intelligence\nfor Fresh Flow Markets",
         24, GRAY, False, "Calibri", PP_ALIGN.CENTER, line_spacing=36)
# divider
add_rect(s, Inches(6.1), Inches(5.4), Inches(1.1), Pt(3), WHITE)
add_text(s, Inches(0), Inches(5.7), SLIDE_W, Inches(0.5),
         "Solving the waste-vs-stockout dilemma across 100+ stores",
         13, GRAY_DARK, False, "Consolas", PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = light_slide(prs)
add_text(s, PAD_L, Inches(0.7), CONTENT_W, Inches(0.3),
         "THE PROBLEM", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(1.0), CONTENT_W, Inches(0.9),
         "Waste vs. Stockouts", 52, BLACK, True)
add_text(s, PAD_L, Inches(1.9), Inches(9.5), Inches(0.7),
         "Fresh Flow Markets faces a dual economic pressure — overordering leads to spoilage "
         "costing DKK millions, while underordering means lost sales at 1.5x the item price.",
         16, GRAY_DARK, False, "Calibri", line_spacing=26)

# Stats row
stats = [("100+", "Stores across Denmark"), ("2M+", "Order items in dataset"), ("1,976", "Unique item-store pairs")]
for i, (val, label) in enumerate(stats):
    x = PAD_L + Inches(i * 3.6)
    # top border
    add_rect(s, x, Inches(3.1), Inches(3.2), Pt(3), BLACK)
    add_text(s, x, Inches(3.3), Inches(3.2), Inches(0.8), val, 48, BLACK, True)
    add_text(s, x, Inches(4.1), Inches(3.2), Inches(0.3), label, 12, GRAY_DARK, False, "Consolas")

# Pain points
pains = [
    ("01", "Gut-feeling ordering with no data-driven approach"),
    ("02", "No adaptation for weather, holidays, or real-time context"),
    ("03", "Static par levels that ignore demand patterns and seasonality"),
]
for i, (num, txt) in enumerate(pains):
    x = PAD_L + Inches(i * 3.6)
    add_rect(s, x, Inches(5.0), Inches(3.2), Inches(1.1), BG_LIGHT)
    add_text(s, x + Inches(0.2), Inches(5.15), Inches(0.5), Inches(0.4), num, 20, BLACK, True)
    add_text(s, x + Inches(0.7), Inches(5.15), Inches(2.3), Inches(0.9), txt, 13, BLACK, False,
             line_spacing=20)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_text(s, PAD_L, Inches(0.6), CONTENT_W, Inches(0.3),
         "OUR SOLUTION", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.9), Inches(8), Inches(1.2),
         "One Platform.\nThree AI Systems.", 52, WHITE, True, line_spacing=62)
add_text(s, PAD_L, Inches(2.3), Inches(9.5), Inches(0.7),
         "FlowPOS integrates a production-grade POS, ML forecasting engine, and LLM intelligence "
         "layer into a single platform that tells managers exactly what to prep — and why.",
         16, GRAY, False, "Calibri", line_spacing=26)

pillars = [
    ("POS Platform",
     "33 pages, 6 roles, 65+ permissions.\nFull checkout, KDS, inventory, loyalty,\ntables, reports — with offline-first PWA."),
    ("ML Forecasting",
     "3-model ensemble: XGBoost + LSTM +\nDeepSeek arbitration. 44 engineered\nfeatures across time, weather, holidays."),
    ("LLM Intelligence",
     "DeepSeek v3.2 with 10 function-calling\ntools. Real-time context signals for\nper-item arbitration."),
]
for i, (title, desc) in enumerate(pillars):
    x = PAD_L + Inches(i * 3.6)
    add_rect(s, x, Inches(3.4), Inches(3.3), Inches(3.2),
             outline=RGBColor(0x33, 0x33, 0x33))
    add_text(s, x + Inches(0.3), Inches(3.7), Inches(2.7), Inches(0.4),
             title, 20, WHITE, True)
    add_text(s, x + Inches(0.3), Inches(4.3), Inches(2.7), Inches(1.8),
             desc, 13, RGBColor(0x88, 0x88, 0x88), False, "Calibri", line_spacing=20)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = light_slide(prs)
add_text(s, PAD_L, Inches(0.5), CONTENT_W, Inches(0.3),
         "ARCHITECTURE", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.8), CONTENT_W, Inches(0.8),
         "System Overview", 48, BLACK, True)

arch_cols = [
    ("FRONTEND", "Next.js 15 + React 19", [
        "33 pages — App Router", "22 Zustand stores",
        "15 tRPC routers (43 endpoints)", "Tailwind v4 + shadcn/ui",
        "English + Arabic (RTL)", "PWA + IndexedDB offline sync"]),
    ("BACKEND / ML", "FastAPI + DuckDB", [
        "18 async API routes", "DuckDB in-memory (2M+ rows)",
        "XGBoost + LSTM ensemble", "44 engineered features",
        "SSE streaming responses", "DeepSeek v3.2 + 10 tools"]),
    ("INFRASTRUCTURE", "Supabase + Railway", [
        "Supabase PostgreSQL + RLS", "Supabase Auth (email, PIN)",
        "Supabase Realtime sync", "Railway containerized deploy",
        "Resend transactional email", "Sentry error monitoring"]),
]
for i, (tag, title, items) in enumerate(arch_cols):
    x = PAD_L + Inches(i * 3.6)
    add_rect(s, x, Inches(1.8), Inches(3.3), Inches(4.8), BLACK)
    add_text(s, x + Inches(0.3), Inches(2.1), Inches(2.7), Inches(0.3),
             tag, 10, GRAY_DARK, False, "Consolas")
    add_text(s, x + Inches(0.3), Inches(2.5), Inches(2.7), Inches(0.5),
             title, 22, WHITE, True)
    lines = [(item, 12, RGBColor(0x88, 0x88, 0x88), False) for item in items]
    add_multiline(s, x + Inches(0.3), Inches(3.3), Inches(2.7), Inches(3.0),
                  lines, font_name="Consolas", line_spacing=22)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: POS PLATFORM
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_text(s, PAD_L, Inches(0.5), CONTENT_W, Inches(0.3),
         "THE PLATFORM", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.8), CONTENT_W, Inches(0.8),
         "Production-Grade POS", 48, WHITE, True)
add_text(s, PAD_L, Inches(1.6), Inches(9), Inches(0.5),
         "A complete restaurant management system — not a prototype. Built for real-world operations from day one.",
         15, RGBColor(0x88, 0x88, 0x88), False, line_spacing=24)

modules = [
    ["POS Terminal", "Kitchen Display", "Inventory Management", "Table Management"],
    ["Employee & Shifts", "Loyalty & Rewards", "Reports & Analytics", "Coupons & Promos"],
    ["Menu Engineering", "Multi-Location", "AI Forecasting", "6 Roles, 65+ Permissions"],
]
for col_i, col in enumerate(modules):
    x = PAD_L + Inches(col_i * 3.6)
    for row_i, mod in enumerate(col):
        y = Inches(2.5) + Inches(row_i * 1.05)
        add_rect(s, x, y, Inches(3.3), Inches(0.8),
                 outline=RGBColor(0x33, 0x33, 0x33))
        add_text(s, x + Inches(0.3), y + Inches(0.2), Inches(2.7), Inches(0.4),
                 mod, 15, WHITE, True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: AI FORECASTING ENGINE
# ═══════════════════════════════════════════════════════════════════════════════
s = light_slide(prs)
add_text(s, PAD_L, Inches(0.5), CONTENT_W, Inches(0.3),
         "AI FORECASTING ENGINE", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.8), CONTENT_W, Inches(0.8),
         "3-Model Ensemble", 48, BLACK, True)
add_text(s, PAD_L, Inches(1.6), Inches(9.5), Inches(0.5),
         "XGBoost captures demand spikes. LSTM learns sequential patterns across 2,604 items. "
         "DeepSeek arbitrates per-item using real-time context.",
         15, GRAY_DARK, False, "Calibri", line_spacing=24)

models = [
    ("01", "XGBoost", "GRADIENT BOOSTING",
     "30% weight in hybrid blend. Trained on\n44 features — captures promotions,\nweather spikes, and calendar effects."),
    ("02", "LSTM (RNN)", "DEEP LEARNING",
     "Unified model for 2,604 items. 2-layer\nLSTM (64→32) with delta-based prediction.\nHandles sparse data and new items."),
    ("03", "DeepSeek LLM", "CONTEXT ARBITRATION",
     "Arbitrates per-item: picks waste-optimized\nor stockout-optimized model based on\nperishability, weather, and events."),
]
for i, (num, title, tag, desc) in enumerate(models):
    x = PAD_L + Inches(i * 3.6)
    add_rect(s, x, Inches(2.5), Inches(3.3), Inches(4.2), BLACK)
    add_text(s, x + Inches(0.3), Inches(2.7), Inches(1), Inches(0.6),
             num, 36, RGBColor(0x33, 0x33, 0x33), True)
    add_text(s, x + Inches(0.3), Inches(3.4), Inches(2.7), Inches(0.5),
             title, 22, WHITE, True)
    add_text(s, x + Inches(0.3), Inches(4.0), Inches(2.7), Inches(1.5),
             desc, 12, RGBColor(0x88, 0x88, 0x88), False, "Calibri", line_spacing=20)
    add_text(s, x + Inches(0.3), Inches(5.8), Inches(2.7), Inches(0.3),
             tag, 10, RGBColor(0x55, 0x55, 0x55), False, "Consolas")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: FEATURES & CONTEXT
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_text(s, PAD_L, Inches(0.5), CONTENT_W, Inches(0.3),
         "44 ENGINEERED FEATURES", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.8), Inches(8), Inches(1.2),
         "Context-Aware\nForecasting", 48, WHITE, True, line_spacing=58)

feat_cards = [
    ("TIME & CALENDAR", "12", "Day of week, month, quarter,\nis_weekend, is_friday, season,\nholidays, bridge days"),
    ("LAG & ROLLING", "12", "1d/7d/14d/28d lags, rolling\nmean & std over 7/14/30 days,\n4-week weekday averages"),
    ("WEATHER", "4", "Max/min temperature,\nprecipitation mm, rain flag\n— via Open-Meteo free API"),
    ("PROMOTIONS & CONTEXT", "16", "Promos, campaigns, cyclical\nencodings, payday proximity,\ndaylight hours, retail events"),
]
for i, (tag, num, desc) in enumerate(feat_cards):
    x = PAD_L + Inches(i * 2.7)
    add_rect(s, x, Inches(2.6), Inches(2.5), Inches(3.0),
             outline=RGBColor(0x33, 0x33, 0x33))
    add_text(s, x + Inches(0.2), Inches(2.8), Inches(2.1), Inches(0.3),
             tag, 9, GRAY_DARK, False, "Consolas")
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(2.1), Inches(0.6),
             num, 40, WHITE, True)
    add_text(s, x + Inches(0.2), Inches(4.0), Inches(2.1), Inches(1.3),
             desc, 11, RGBColor(0x88, 0x88, 0x88), False, "Calibri", line_spacing=18)

# Signal badges
signals = ["Weather", "Holidays", "Daylight", "Payday", "Retail Events"]
add_text(s, PAD_L, Inches(6.2), Inches(1.5), Inches(0.3),
         "REAL-TIME SIGNALS", 9, RGBColor(0x55, 0x55, 0x55), False, "Consolas")
for i, sig in enumerate(signals):
    x = PAD_L + Inches(1.8) + Inches(i * 1.6)
    add_rect(s, x, Inches(6.1), Inches(1.4), Inches(0.4),
             outline=RGBColor(0x33, 0x33, 0x33))
    add_text(s, x, Inches(6.15), Inches(1.4), Inches(0.3),
             sig, 11, WHITE, False, "Consolas", PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: LLM INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════════
s = light_slide(prs)
add_text(s, PAD_L, Inches(0.5), CONTENT_W, Inches(0.3),
         "LLM INTELLIGENCE LAYER", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.8), CONTENT_W, Inches(0.8),
         "DeepSeek + 10 Tools", 48, BLACK, True)
add_text(s, PAD_L, Inches(1.6), Inches(9.5), Inches(0.6),
         "Manager asks a question → LLM calls tools → fetches context → queries models → "
         "arbitrates per item → returns structured prep recommendations with risk flags.",
         14, GRAY_DARK, False, "Calibri", line_spacing=22)

tools_left = [
    "query_inventory", "get_sales_history", "get_forecast",
    "get_dual_forecast", "get_context_signals"
]
tools_right = [
    "get_low_stock", "get_expiring_items", "get_top_sellers",
    "get_bill_of_materials", "run_sql"
]
for i, tool in enumerate(tools_left):
    y = Inches(2.6) + Inches(i * 0.55)
    add_rect(s, PAD_L, y, Inches(5.0), Inches(0.45), BG_LIGHT)
    add_text(s, PAD_L + Inches(0.2), y + Inches(0.08), Inches(4.5), Inches(0.3),
             tool, 12, BLACK, True, "Consolas")
for i, tool in enumerate(tools_right):
    y = Inches(2.6) + Inches(i * 0.55)
    add_rect(s, PAD_L + Inches(5.4), y, Inches(5.0), Inches(0.45), BG_LIGHT)
    add_text(s, PAD_L + Inches(5.6), y + Inches(0.08), Inches(4.5), Inches(0.3),
             tool, 12, BLACK, True, "Consolas")

# AI Personas
add_text(s, PAD_L, Inches(5.6), Inches(1.3), Inches(0.3),
         "AI PERSONAS", 9, GRAY_DARK, False, "Consolas")
personas = ["Chat Assistant", "Inventory Analyst", "Anomaly Explainer",
            "Promo Strategist", "Scenario Simulator", "Daily Briefing"]
for i, p in enumerate(personas):
    x = PAD_L + Inches(1.5) + Inches(i * 1.7)
    add_rect(s, x, Inches(5.5), Inches(1.55), Inches(0.45), BLACK)
    add_text(s, x, Inches(5.57), Inches(1.55), Inches(0.3),
             p, 10, WHITE, True, "Consolas", PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: RESULTS & IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_text(s, PAD_L, Inches(0.5), CONTENT_W, Inches(0.3),
         "RESULTS & IMPACT", 12, GRAY_DARK, False, "Consolas")
add_text(s, PAD_L, Inches(0.8), CONTENT_W, Inches(1.0),
         "27.8% Cost Reduction", 56, WHITE, True)
add_text(s, PAD_L, Inches(1.8), Inches(9.5), Inches(0.5),
         "Measured against the worst baseline over a 93-day test period across 101 stores. "
         "Business-cost metric: waste DKK + 1.5x stockout DKK.",
         15, RGBColor(0x88, 0x88, 0x88), False, "Calibri", line_spacing=24)

results = [
    ("67.7%", "Forecast Accuracy (WMAPE)"),
    ("38.1%", "Balanced Forecast Days"),
    ("15.5M", "DKK Total Cost (optimized)"),
    ("93", "Days Test Period"),
]
for i, (val, label) in enumerate(results):
    x = PAD_L + Inches(i * 2.7)
    add_rect(s, x, Inches(2.9), Inches(2.5), Pt(3), WHITE)
    add_text(s, x, Inches(3.1), Inches(2.5), Inches(0.7), val, 42, WHITE, True)
    add_text(s, x, Inches(3.85), Inches(2.5), Inches(0.3), label, 11, GRAY_DARK, False, "Consolas")

# Model comparison
add_text(s, PAD_L, Inches(4.8), Inches(1.8), Inches(0.3),
         "MODEL COMPARISON", 9, RGBColor(0x55, 0x55, 0x55), False, "Consolas")
model_comp = [
    ("Balanced", "30% XGB + 70% MA7"),
    ("Waste-Optimized", "85% of balanced (15% shrink)"),
    ("LSTM Unified", "2,604 items, 1.04 MAE"),
    ("32 Configs", "Grid search evaluated"),
]
for i, (title, sub) in enumerate(model_comp):
    x = PAD_L + Inches(2.2) + Inches(i * 2.3)
    add_text(s, x, Inches(4.7), Inches(2.0), Inches(0.3), title, 16, WHITE, True)
    add_text(s, x, Inches(5.1), Inches(2.0), Inches(0.3), sub, 11, GRAY_DARK, False, "Consolas")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════════════
s = light_slide(prs)
add_text(s, Inches(0), Inches(1.5), SLIDE_W, Inches(0.4),
         "LIVE DEMO", 13, GRAY_DARK, False, "Consolas", PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(2.0), SLIDE_W, Inches(1.5),
         "Let's See It\nin Action", 64, BLACK, True, "Calibri", PP_ALIGN.CENTER, line_spacing=76)
# divider
add_rect(s, Inches(6.1), Inches(3.9), Inches(1.1), Pt(3), BLACK)
# Links
add_rect(s, Inches(3.2), Inches(4.4), Inches(6.9), Inches(0.6), BLACK)
add_text(s, Inches(3.4), Inches(4.48), Inches(6.5), Inches(0.4),
         "POS Frontend — pos-frontend-production-56bb.up.railway.app",
         13, WHITE, True, "Consolas", PP_ALIGN.CENTER)
add_rect(s, Inches(3.2), Inches(5.15), Inches(6.9), Inches(0.6),
         outline=BLACK)
add_text(s, Inches(3.4), Inches(5.23), Inches(6.5), Inches(0.4),
         "Forecasting API — hopeful-elegance-production-c09a.up.railway.app",
         13, BLACK, True, "Consolas", PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.1), SLIDE_W, Inches(0.4),
         "Both services are live and deployed on Railway",
         12, GRAY, False, "Consolas", PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11: THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_slide(prs)
add_text(s, Inches(0), Inches(1.5), SLIDE_W, Inches(0.4),
         "DIH x AUC HACKATHON", 13, GRAY_DARK, False, "Consolas", PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(2.2), SLIDE_W, Inches(1.5),
         "Thank You", 80, WHITE, True, "Calibri", PP_ALIGN.CENTER)
# divider
add_rect(s, Inches(6.1), Inches(3.8), Inches(1.1), Pt(3), WHITE)
add_text(s, Inches(0), Inches(4.1), SLIDE_W, Inches(0.5),
         "FlowPOS — AI-Powered Demand Intelligence",
         20, RGBColor(0x88, 0x88, 0x88), False, "Calibri", PP_ALIGN.CENTER)

# Stats row
ty_stats = [("33", "Pages"), ("44", "Features"), ("10", "LLM Tools"), ("3", "AI Models")]
total_w = len(ty_stats) * 1.5 + (len(ty_stats) - 1) * 0.6
start_x = (13.333 - total_w) / 2
for i, (val, label) in enumerate(ty_stats):
    x = Inches(start_x + i * 2.1)
    add_text(s, x, Inches(4.9), Inches(1.5), Inches(0.5),
             val, 32, WHITE, True, "Calibri", PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.4), Inches(1.5), Inches(0.3),
             label, 11, GRAY_DARK, False, "Consolas", PP_ALIGN.CENTER)
    if i < len(ty_stats) - 1:
        add_text(s, Inches(start_x + i * 2.1 + 1.6), Inches(5.0), Inches(0.3), Inches(0.4),
                 "/", 20, RGBColor(0x33, 0x33, 0x33), False, "Calibri", PP_ALIGN.CENTER)

add_text(s, Inches(0), Inches(6.1), SLIDE_W, Inches(0.4),
         "AUC Team", 14, RGBColor(0x55, 0x55, 0x55), True, "Consolas", PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
output = "/home/yahyahammoudeh/Documents/Dih/DIH-X-AUC-Hackathon/FlowPOS-Hackathon-Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
